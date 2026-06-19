#!/usr/bin/env python3
"""Convert Cashfree Relay pitch deck HTML → .pptx (Google Slides compatible)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Emu(9144000)   # 10 inches  (16:9 Google Slides)
prs.slide_height = Emu(5143500)   # 5.625 inches

W = prs.slide_width
H = prs.slide_height
BL = prs.slide_layouts[6]  # blank layout

# ── palette ──────────────────────────────────────────────────────────
GREEN      = RGBColor(0x00, 0xB4, 0x88)
GREEN_INK  = RGBColor(0x00, 0x7A, 0x55)
GREEN_SOFT = RGBColor(0x2D, 0xCC, 0xA0)
GREEN_TINT = RGBColor(0xE8, 0xF8, 0xF3)
INK        = RGBColor(0x0F, 0x1F, 0x1B)
INK_MID    = RGBColor(0x44, 0x52, 0x4C)
MUTED      = RGBColor(0x5E, 0x6A, 0x66)
AMBER      = RGBColor(0x92, 0x40, 0x0E)
RED        = RGBColor(0xC0, 0x39, 0x2B)
RED_LIGHT  = RGBColor(0xF5, 0xEC, 0xEB)
PAPER      = RGBColor(0xFB, 0xFB, 0xFB)
DARK       = RGBColor(0x0F, 0x1F, 0x1B)
OFF_WHITE  = RGBColor(0xF4, 0xF1, 0xEA)
LINE       = RGBColor(0xD6, 0xDD, 0xDB)
DARK_LINE  = RGBColor(0x1E, 0x31, 0x2B)

FONT = 'Calibri'  # closest available to Google Sans

# ── helpers ───────────────────────────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(BL)

def set_bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def box(slide, x, y, w, h, fill=None, line_color=None, line_w=Pt(0.75)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def tb(slide, text, x, y, w, h,
       size=Pt(11), bold=False, italic=False, color=INK,
       align=PP_ALIGN.LEFT, wrap=True, spacing_after=0):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing_after:
        p.space_after = Pt(spacing_after)
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = size
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return shape

def para(tf, text, size=Pt(11), bold=False, italic=False, color=INK,
         align=PP_ALIGN.LEFT, space_before=0, space_after=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before: p.space_before = Pt(space_before)
    if space_after:  p.space_after  = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = size
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return p

def chrome_bar(slide, num, dark=False):
    """Top brand mark + bottom slide number"""
    txt_c = OFF_WHITE if dark else MUTED
    num_c = OFF_WHITE if dark else INK
    dot_c = GREEN_SOFT if dark else GREEN
    # green dot
    box(slide, Inches(0.28), Inches(0.15), Inches(0.075), Inches(0.075), fill=dot_c)
    # brand
    tb(slide, 'Cashfree · Relay', Inches(0.40), Inches(0.10), Inches(2.5), Inches(0.28),
       size=Pt(7.5), bold=True, color=txt_c)
    # slide number
    tb(slide, num, Inches(9.5), H - Inches(0.32), Inches(0.5), Inches(0.25),
       size=Pt(8), bold=True, color=num_c, align=PP_ALIGN.RIGHT)
    # top rule
    box(slide, Inches(0.28), Inches(0.36), W - Inches(0.56), Emu(6350),
        fill=dot_c if dark else LINE)

def eyebrow_label(slide, text, x, y, w, color=MUTED):
    box(slide, x, y + Pt(5), Inches(0.2), Emu(9525), fill=GREEN)
    tb(slide, text.upper(), x + Inches(0.28), y, w, Inches(0.22),
       size=Pt(7.5), bold=True, color=color)

def section_card(slide, x, y, w, h, header, header_bg, header_color, body_text,
                 border_color=LINE, footer=None, footer_bg=None, footer_color=INK_MID):
    """Generic card: header strip + body + optional footer"""
    hh = Inches(0.32)
    box(slide, x, y, w, hh, fill=header_bg, line_color=border_color)
    tb(slide, header.upper(), x + Inches(0.1), y + Pt(4), w - Inches(0.2), hh - Pt(4),
       size=Pt(7), bold=True, color=header_color)
    bh = h - hh - (Inches(0.28) if footer else 0)
    box(slide, x, y + hh, w, bh, fill=PAPER, line_color=border_color, line_w=Pt(0.5))
    body_tb = slide.shapes.add_textbox(x + Inches(0.1), y + hh + Pt(6),
                                        w - Inches(0.2), bh - Pt(12))
    body_tb.text_frame.word_wrap = True
    p0 = body_tb.text_frame.paragraphs[0]
    r = p0.add_run()
    r.text = body_text
    r.font.name = FONT
    r.font.size = Pt(9.5)
    r.font.color.rgb = INK_MID
    if footer:
        fh = Inches(0.28)
        box(slide, x, y + hh + bh, w, fh,
            fill=footer_bg or PAPER, line_color=border_color, line_w=Pt(0.5))
        tb(slide, footer, x + Inches(0.1), y + hh + bh + Pt(4),
           w - Inches(0.2), fh, size=Pt(8.5), bold=True, color=footer_color)

def quote_block(slide, quote_text, source, x, y, w, h):
    box(slide, x, y, Emu(15875), h, fill=AMBER)  # amber left border ~2px
    box(slide, x + Emu(15875), y, w - Emu(15875), h, fill=RGBColor(0xF9, 0xF4, 0xF2))
    quote_tb = slide.shapes.add_textbox(x + Inches(0.15), y + Pt(6),
                                          w - Inches(0.2), h - Inches(0.35))
    quote_tb.text_frame.word_wrap = True
    p = quote_tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = quote_text
    r.font.name = FONT
    r.font.size = Pt(9)
    r.font.italic = True
    r.font.color.rgb = INK_MID
    tb(slide, source.upper(), x + Inches(0.15), y + h - Inches(0.22),
       w - Inches(0.2), Inches(0.2), size=Pt(7), bold=True, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 00 · TITLE
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, PAPER)
chrome_bar(s, '00')

# green accent panel (right)
box(s, W - Inches(3.8), 0, Inches(3.8), H, fill=RGBColor(0xF5, 0xFD, 0xFA))
box(s, W - Inches(3.8), 0, Emu(9525), H, fill=RGBColor(0xD0, 0xF0, 0xE8))

# headline
hl = s.shapes.add_textbox(Inches(0.55), Inches(1.4), Inches(6.2), Inches(2.2))
hl.text_frame.word_wrap = True
p1 = hl.text_frame.paragraphs[0]
r1 = p1.add_run(); r1.text = 'For years, software helped you '
r1.font.name = FONT; r1.font.size = Pt(38); r1.font.color.rgb = INK; r1.font.bold = False
p2 = hl.text_frame.add_paragraph()
r2 = p2.add_run(); r2.text = 'see the problem.'
r2.font.name = FONT; r2.font.size = Pt(38); r2.font.color.rgb = GREEN; r2.font.italic = True

# lede
tb(s, 'This is software that fixes it. Cashfree Relay puts specialized AI agents on your payment events. Automatically. Without you.',
   Inches(0.55), Inches(3.85), Inches(5.8), Inches(0.9),
   size=Pt(13), italic=True, color=INK_MID)

# meta
tb(s, 'Cashfree Payments  ·  Early Beta  ·  Free to activate',
   Inches(0.55), Inches(4.85), Inches(5), Inches(0.3),
   size=Pt(8.5), bold=True, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 01 · DARK DIVIDER — YOUR MORNING
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, DARK)
chrome_bar(s, '01 / 10', dark=True)

eyebrow_label(s, 'Before we talk about the product', Inches(0.55), Inches(1.5),
              Inches(6), color=RGBColor(0x88, 0x9E, 0x98))

hl = s.shapes.add_textbox(Inches(0.55), Inches(1.9), Inches(9), Inches(2.5))
hl.text_frame.word_wrap = True
p1 = hl.text_frame.paragraphs[0]
r1 = p1.add_run(); r1.text = 'Walk us through '
r1.font.name = FONT; r1.font.size = Pt(52); r1.font.color.rgb = OFF_WHITE

p2 = hl.text_frame.add_paragraph()
r2 = p2.add_run(); r2.text = 'your Monday morning.'
r2.font.name = FONT; r2.font.size = Pt(52); r2.font.color.rgb = GREEN_SOFT; r2.font.italic = True


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 02 · THE MORNING
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, PAPER)
chrome_bar(s, '02 / 10')

eyebrow_label(s, 'D2C founder · ₹30L–5Cr monthly GMV · 40%+ COD', Inches(0.55), Inches(0.5), Inches(7))

hl = s.shapes.add_textbox(Inches(0.55), Inches(0.8), Inches(9), Inches(0.65))
hl.text_frame.word_wrap = True
p = hl.text_frame.paragraphs[0]
r1 = p.add_run(); r1.text = 'You open your phone at 8am. '
r1.font.name = FONT; r1.font.size = Pt(20); r1.font.color.rgb = INK
r2 = p.add_run(); r2.text = 'The ops work is already waiting.'
r2.font.name = FONT; r2.font.size = Pt(20); r2.font.color.rgb = INK; r2.font.italic = True

# Left col — time steps
def step_block(slide, label, body, x, y, w, h):
    box(slide, x, y, Emu(12700), h, fill=RED)
    box(slide, x + Emu(12700), y, w - Emu(12700), h, fill=RED_LIGHT)
    tb(slide, label.upper(), x + Inches(0.18), y + Pt(4), w - Inches(0.2), Inches(0.2),
       size=Pt(7), bold=True, color=RED)
    tb(slide, body, x + Inches(0.18), y + Inches(0.25), w - Inches(0.25), h - Inches(0.3),
       size=Pt(9), color=INK_MID, wrap=True)

step_block(s, '8:05am · COD confirmations',
           '47 COD orders came in overnight. You open WhatsApp and start sending confirmation messages one by one. This takes until 9:30am. Three people didn\'t pick up.',
           Inches(0.45), Inches(1.6), Inches(4.4), Inches(1.05))
step_block(s, '9:45am · Failed payments',
           '14 payments failed yesterday. You export the list, open a spreadsheet, start writing recovery messages. You\'ve sent 6 when a supplier calls. You\'ll finish the rest after lunch. You probably won\'t.',
           Inches(0.45), Inches(2.72), Inches(4.4), Inches(1.05))
step_block(s, 'After lunch · Settlement recon',
           'Settlement ran on Friday. Finance is asking if the numbers match. You export the CSV, open Sheets, start matching rows. This takes the rest of the afternoon. You\'ve been doing this every week for 14 months.',
           Inches(0.45), Inches(3.84), Inches(4.4), Inches(1.05))

# Right col — quotes
qy = Inches(1.6)
qh = Inches(1.05)
qgap = Inches(0.08)
quote_block(s,
    '"Last month alone, 23% of COD orders were fake or undelivered. Average loss per order around ₹350. We\'ve tried OTP verification, restricting postal codes, charging extra for COD. Nothing is really helping."',
    'r/IndiaStartups · Nov 2025 · 78 upvotes',
    Inches(5.1), qy, Inches(4.7), qh)
quote_block(s,
    '"The customer dropped off. I wanted a nudge to go automatically. I was sending it manually."',
    'Merchant founder · Ohsou · ₹20L GMV/month · Cashfree demo call',
    Inches(5.1), qy + qh + qgap, Inches(4.7), Inches(0.82))
quote_block(s,
    '"Version 1 was embarrassing. The final system took 2 months — 44 nodes, 2 AI agents, 8 conversation stages. Just for WhatsApp replies."',
    'r/n8n · March 2026 · 356 upvotes · Developer, Jaipur',
    Inches(5.1), qy + qh + qgap + Inches(0.82) + qgap, Inches(4.7), Inches(0.88))

# dark bar
box(s, Inches(5.1), Inches(4.5), Inches(4.7), Inches(0.42), fill=DARK)
tb(s, '271 posts scraped', Inches(5.2), Inches(4.55), Inches(2), Inches(0.32),
   size=Pt(7.5), bold=True, color=RGBColor(0x55, 0x6A, 0x60))
tb(s, 'Zero merchants said "we use X and it works."',
   Inches(7.0), Inches(4.55), Inches(2.7), Inches(0.32),
   size=Pt(9.5), color=GREEN_SOFT, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 03 · THE COST TABLE
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, PAPER)
chrome_bar(s, '03 / 10')

eyebrow_label(s, "You're already paying for this problem", Inches(0.55), Inches(0.5), Inches(6))

hl = s.shapes.add_textbox(Inches(0.55), Inches(0.8), Inches(9), Inches(0.65))
p = hl.text_frame.paragraphs[0]
r1 = p.add_run(); r1.text = '₹40,000 to ₹80,000 every month. '
r1.font.name = FONT; r1.font.size = Pt(22); r1.font.color.rgb = INK
r2 = p.add_run(); r2.text = 'And it still breaks on weekends.'
r2.font.name = FONT; r2.font.size = Pt(22); r2.font.color.rgb = INK; r2.font.italic = True

# Table
rows = [
    ('WhatsApp COD confirmations. Sent one by one, every morning.', '₹8,000 – ₹18,000 / mo'),
    ('Connecting Cashfree events to Sheets, CRM, Slack. Done manually.', '₹4,000 – ₹12,000 / mo'),
    ('COD RTO ops. Calling, flagging, chasing before dispatch.', '₹20,000 – ₹35,000 / mo'),
    ('Failed payment follow-ups. Export, message, track manually.', '2 hrs every day'),
    ('Settlement reconciliation. CSV export, Sheets matching, finance review.', '40 – 160 hrs / month'),
]
tx = Inches(0.45)
ty = Inches(1.6)
tw = Inches(9.1)
rh = Inches(0.5)
col1w = Inches(6.8)
col2w = tw - col1w

# header row
box(s, tx, ty, tw, Inches(0.32), fill=RGBColor(0xED, 0xF0, 0xEF), line_color=INK)
tb(s, 'WHAT YOU\'RE DOING MANUALLY', tx + Inches(0.1), ty + Pt(4),
   col1w - Inches(0.1), Inches(0.28), size=Pt(7.5), bold=True, color=MUTED)
tb(s, 'WHAT IT COSTS', tx + col1w + Inches(0.1), ty + Pt(4),
   col2w - Inches(0.1), Inches(0.28), size=Pt(7.5), bold=True, color=MUTED)

for i, (task, cost) in enumerate(rows):
    ry = ty + Inches(0.32) + i * rh
    bg_fill = PAPER if i % 2 == 0 else RGBColor(0xF5, 0xF7, 0xF6)
    box(s, tx, ry, col1w, rh, fill=bg_fill, line_color=LINE, line_w=Pt(0.5))
    box(s, tx + col1w, ry, col2w, rh, fill=bg_fill, line_color=LINE, line_w=Pt(0.5))
    tb(s, task, tx + Inches(0.1), ry + Pt(4), col1w - Inches(0.2), rh - Pt(8),
       size=Pt(9.5), bold=True, color=INK, wrap=True)
    tb(s, cost, tx + col1w + Inches(0.1), ry + Pt(10), col2w - Inches(0.1), rh - Pt(8),
       size=Pt(11), bold=True, color=AMBER)

# quote + dark bar
qy2 = ty + Inches(0.32) + len(rows) * rh + Inches(0.1)
qh2 = Inches(0.5)
quote_block(s,
    '"I was paying ₹6,000/month for automation tools and a freelance developer to maintain three workflows. Still had to check everything myself."',
    'Early Beta merchant · D2C apparel brand',
    tx, qy2, Inches(4.4), qh2)
box(s, tx + Inches(4.6), qy2, Inches(4.5), qh2, fill=DARK)
tb(s, 'The problem isn\'t the cost', tx + Inches(4.7), qy2 + Pt(4), Inches(2.2), Inches(0.22),
   size=Pt(7.5), bold=True, color=RGBColor(0x55, 0x6A, 0x60))
tb(s, 'It\'s that you\'re spending this money and still doing the work.',
   tx + Inches(7.0), qy2 + Pt(4), Inches(2), qh2,
   size=Pt(9.5), color=GREEN_SOFT, align=PP_ALIGN.RIGHT, wrap=True)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 04 · DARK DIVIDER — THE SHIFT
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, DARK)
chrome_bar(s, '04 / 10', dark=True)

eyebrow_label(s, 'The shift', Inches(0.55), Inches(1.3),
              Inches(4), color=RGBColor(0x88, 0x9E, 0x98))

hl = s.shapes.add_textbox(Inches(0.55), Inches(1.7), Inches(9), Inches(1.9))
hl.text_frame.word_wrap = True
p1 = hl.text_frame.paragraphs[0]
r1 = p1.add_run(); r1.text = 'Not another dashboard.'
r1.font.name = FONT; r1.font.size = Pt(48); r1.font.color.rgb = OFF_WHITE
p2 = hl.text_frame.add_paragraph()
r2 = p2.add_run(); r2.text = 'A team that does the work.'
r2.font.name = FONT; r2.font.size = Pt(48); r2.font.color.rgb = GREEN_SOFT; r2.font.italic = True

tb(s, '"For years, software had businesses see what is happening — dashboards, reports, alerts. But the work, you still had to do it. What you see today is different."',
   Inches(0.55), Inches(3.75), Inches(7.5), Inches(0.9),
   size=Pt(13), italic=True, color=RGBColor(0x88, 0x9E, 0x98), wrap=True)

tb(s, 'THE SHIFT THE ENTIRE INDUSTRY IS MAKING',
   Inches(0.55), Inches(4.72), Inches(6), Inches(0.28),
   size=Pt(7), bold=True, color=RGBColor(0x44, 0x56, 0x50))


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 05 · WHAT RELAY IS
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, PAPER)
chrome_bar(s, '05 / 10')

eyebrow_label(s, 'One sentence', Inches(0.55), Inches(0.5), Inches(4))

hl = s.shapes.add_textbox(Inches(0.55), Inches(0.8), Inches(9), Inches(0.75))
hl.text_frame.word_wrap = True
p = hl.text_frame.paragraphs[0]
r1 = p.add_run()
r1.text = 'Specialized AI agents for every post-payment problem. '
r1.font.name = FONT; r1.font.size = Pt(20); r1.font.color.rgb = INK
r2 = p.add_run()
r2.text = 'Activate once. They run on every event after that.'
r2.font.name = FONT; r2.font.size = Pt(20); r2.font.color.rgb = INK; r2.font.italic = True

cards = [
    {
        'header': 'Event-triggered or scheduled',
        'header_bg': GREEN_TINT, 'header_color': GREEN_INK,
        'border': GREEN,
        'body': 'Trigger on a Cashfree payment event (COD placed, payment failed, refund initiated, dispute raised) and the AI agent acts within minutes. Or run on a schedule: daily recon, weekly reports, recurring recovery sweeps. Both are built in.',
        'footer': 'COD order in at 2pm · confirmation out by 2:05pm',
        'footer_bg': RGBColor(0xD6, 0xF5, 0xEB), 'footer_color': GREEN_INK,
    },
    {
        'header': 'No developer needed',
        'header_bg': RGBColor(0xF0, 0xF2, 0xF1), 'header_color': MUTED,
        'border': LINE,
        'body': 'Activate from inside the Cashfree dashboard. Set your guardrails once: what to send, when to escalate, when to flag you. The AI agent follows those rules on every qualifying event. No webhook setup. No API keys to manage.',
        'footer': 'You set the rules once. The AI agent runs forever.',
        'footer_bg': RGBColor(0xF7, 0xF8, 0xF8), 'footer_color': INK_MID,
    },
    {
        'header': 'Data stays in Cashfree',
        'header_bg': RGBColor(0xF0, 0xF2, 0xF1), 'header_color': MUTED,
        'border': LINE,
        'body': 'Relay runs inside Cashfree\'s PCI DSS Level 1 and SOC infrastructure. Google, Slack, and WhatsApp only receive the specific fields you send them in each step. Nothing else leaves.',
        'footer': 'You control the AI agents. Cashfree runs them.',
        'footer_bg': RGBColor(0xF7, 0xF8, 0xF8), 'footer_color': INK_MID,
    },
]
cw = Inches(2.9)
cgap = Inches(0.2)
cx_start = Inches(0.45)
cy = Inches(1.7)
ch = Inches(3.1)

for i, c in enumerate(cards):
    cx = cx_start + i * (cw + cgap)
    section_card(s, cx, cy, cw, ch,
                 c['header'], c['header_bg'], c['header_color'],
                 c['body'], border_color=c['border'],
                 footer=c['footer'], footer_bg=c['footer_bg'], footer_color=c['footer_color'])

# insight
box(s, Inches(0.45), Inches(4.95), Inches(9.1), Inches(0.38),
    fill=RGBColor(0xF0, 0xF9, 0xF5))
tb(s, '=> Most tools need your payment data passed to them. Relay already has it. That is the entire difference.',
   Inches(0.55), Inches(5.0), Inches(9), Inches(0.3),
   size=Pt(9.5), italic=True, color=INK_MID)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 06 · CONNECTORS
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, PAPER)
chrome_bar(s, '06 / 10')

eyebrow_label(s, 'Trigger from Cashfree. Act everywhere else.', Inches(0.55), Inches(0.5), Inches(7))
tb(s, 'Your payment event happens.', Inches(0.55), Inches(0.82), Inches(5.5), Inches(0.45),
   size=Pt(22), color=INK)
tb(s, 'Relay takes it from there.', Inches(0.55), Inches(1.22), Inches(5.5), Inches(0.45),
   size=Pt(22), italic=True, color=INK)

# Column positions
col1x = Inches(0.45); col1w = Inches(2.5)
col2x = Inches(3.25); col2w = Inches(3.8)
col3x = Inches(7.35); col3w = Inches(2.5)
col_y = Inches(1.8)
col_h = Inches(3.5)

# Col 1 - Trigger
tb(s, 'TRIGGER', col1x, col_y, col1w, Inches(0.25), size=Pt(7.5), bold=True, color=GREEN_INK)
box(s, col1x, col_y + Inches(0.28), col1w, Inches(1.5), fill=GREEN_TINT, line_color=GREEN, line_w=Pt(1.2))
tb(s, 'Cashfree Payments', col1x + Inches(0.12), col_y + Inches(0.36), col1w - Inches(0.2), Inches(0.28),
   size=Pt(11), bold=True, color=INK)
tb(s, '11 actions: payments, refunds, disputes, settlement, payment links, cashgrams.',
   col1x + Inches(0.12), col_y + Inches(0.65), col1w - Inches(0.2), Inches(0.7),
   size=Pt(9), color=INK_MID, wrap=True)
tb(s, 'Your existing Cashfree account. No extra setup.',
   col1x + Inches(0.12), col_y + Inches(1.35), col1w - Inches(0.2), Inches(0.35),
   size=Pt(8.5), bold=True, color=GREEN_INK, wrap=True)

box(s, col1x, col_y + Inches(1.9), col1w, Inches(1.0), fill=PAPER, line_color=LINE, line_w=Pt(0.75))
tb(s, 'Schedule', col1x + Inches(0.12), col_y + Inches(1.98), col1w - Inches(0.2), Inches(0.28),
   size=Pt(11), bold=True, color=INK)
tb(s, 'Daily, weekly, or custom cadence. Recon sweeps, recovery batches, periodic reports.',
   col1x + Inches(0.12), col_y + Inches(2.28), col1w - Inches(0.2), Inches(0.58),
   size=Pt(9), color=INK_MID, wrap=True)

# Divider
box(s, col1x + col1w + Inches(0.12), col_y, Emu(9525), col_h, fill=LINE)

# Col 2 - Act via
tb(s, 'ACT VIA', col2x, col_y, col2w, Inches(0.25), size=Pt(7.5), bold=True, color=MUTED)
act_items = [
    ('WhatsApp Business', '2 actions'),
    ('Google Sheets',     '1 action'),
    ('Slack',             '4 actions'),
    ('Gmail',             '2 actions'),
    ('Google Calendar',   '2 actions'),
    ('Calendly',          '1 action'),
    ('HTTP',              '1 action · call any endpoint'),
]
iw = (col2w - Inches(0.1)) / 2
for idx, (name, actions) in enumerate(act_items):
    irow = idx // 2; icol = idx % 2
    if name == 'HTTP':
        ix = col2x; iw2 = col2w
    else:
        ix = col2x + icol * (iw + Inches(0.08)); iw2 = iw
    iy = col_y + Inches(0.28) + irow * Inches(0.72)
    box(s, ix, iy, iw2, Inches(0.65), fill=PAPER, line_color=LINE, line_w=Pt(0.5))
    tb(s, name, ix + Inches(0.1), iy + Pt(5), iw2 - Inches(0.15), Inches(0.28),
       size=Pt(10), bold=True, color=INK)
    tb(s, actions, ix + Inches(0.1), iy + Inches(0.33), iw2 - Inches(0.15), Inches(0.25),
       size=Pt(8), color=MUTED)

# Divider
box(s, col3x - Inches(0.12), col_y, Emu(9525), col_h, fill=LINE)

# Col 3 - AI inside
tb(s, 'AI INSIDE THE AGENT', col3x, col_y, col3w, Inches(0.25), size=Pt(7.5), bold=True, color=MUTED)
box(s, col3x, col_y + Inches(0.28), col3w, Inches(1.2), fill=PAPER, line_color=LINE, line_w=Pt(0.75))
tb(s, 'Google Gemini', col3x + Inches(0.12), col_y + Inches(0.36), col3w - Inches(0.2), Inches(0.28),
   size=Pt(11), bold=True, color=INK)
tb(s, '6 actions. Add an AI decision step inside any agent: classify an order, draft a reply, score a risk.',
   col3x + Inches(0.12), col_y + Inches(0.66), col3w - Inches(0.2), Inches(0.75),
   size=Pt(9), color=INK_MID, wrap=True)

box(s, col3x, col_y + Inches(1.6), col3w, Inches(1.2), fill=PAPER, line_color=LINE, line_w=Pt(0.75))
tb(s, 'OpenAI', col3x + Inches(0.12), col_y + Inches(1.68), col3w - Inches(0.2), Inches(0.28),
   size=Pt(11), bold=True, color=INK)
tb(s, '9 actions. Use GPT models for reasoning, summarisation, or classification steps inside your agent.',
   col3x + Inches(0.12), col_y + Inches(1.98), col3w - Inches(0.2), Inches(0.75),
   size=Pt(9), color=INK_MID, wrap=True)

box(s, col3x, col_y + Inches(2.92), col3w, Inches(0.45), fill=DARK)
tb(s, 'AI is a step inside your agent. Not a separate subscription.',
   col3x + Inches(0.1), col_y + Inches(3.0), col3w - Inches(0.15), Inches(0.35),
   size=Pt(9), color=GREEN_SOFT, wrap=True)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 07 · FIVE AI AGENTS
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, PAPER)
chrome_bar(s, '07 / 10')

eyebrow_label(s, 'Pick the one that hurts most right now', Inches(0.55), Inches(0.5), Inches(6))

hl = s.shapes.add_textbox(Inches(0.55), Inches(0.82), Inches(9), Inches(0.5))
p = hl.text_frame.paragraphs[0]
r1 = p.add_run(); r1.text = 'Five things you do manually today. '
r1.font.name = FONT; r1.font.size = Pt(20); r1.font.color.rgb = INK
r2 = p.add_run(); r2.text = 'Five AI agents that do them instead.'
r2.font.name = FONT; r2.font.size = Pt(20); r2.font.color.rgb = INK; r2.font.bold = True

agents = [
    {
        'tag': 'COD', 'tag_color': GREEN, 'tag_bg': GREEN_TINT,
        'note': 'Most activated', 'note_color': GREEN_INK,
        'border': GREEN,
        'title': 'COD Confirmation AI Agent',
        'body': 'Every COD order confirmed within 5 minutes. No reply in 4 hours? A nudge goes out. Still nothing? Ops is flagged before dispatch. Runs at 2am. Runs on Sunday. You don\'t send a single message.',
        'stat': 'RTO: 28–32% → 16–20%',
        'stat_note': 'Within 30 days · CampaignHQ 2026',
    },
    {
        'tag': 'Recovery', 'tag_color': INK, 'tag_bg': RGBColor(0xED, 0xEF, 0xEE),
        'border': LINE,
        'title': 'Failed Payment Recovery',
        'body': 'WhatsApp within 15 minutes. Gmail backup at 2 hours. Slack alert to ops at end of day if still unpaid. Covers one-time payments, COD, payment links, checkout drops.',
        'stat': '3–8% of drops recovered',
        'stat_note': 'Early Beta merchants · no ops involvement',
    },
    {
        'tag': 'Cart', 'tag_color': INK, 'tag_bg': RGBColor(0xED, 0xEF, 0xEE),
        'border': LINE,
        'title': 'Abandoned Cart Recovery',
        'body': 'Customer drops at checkout. WhatsApp within 10 minutes. No reply at 2 hours. Follow-up with an offer goes out. Triggered from the payment drop event, not browsing.',
        'stat': 'Timed to highest intent',
        'stat_note': 'Triggered at payment drop · not browsing',
    },
    {
        'tag': 'Recon', 'tag_color': INK, 'tag_bg': RGBColor(0xED, 0xEF, 0xEE),
        'border': LINE,
        'title': 'Settlement Recon AI Agent',
        'body': 'Runs daily, weekly, or after every settlement. Pulls every row, appends to Google Sheets, flags discrepancies, notifies finance. A live tab your finance team can act on.',
        'stat': '40–160 hrs / month → 0',
        'stat_note': 'Kani survey 2025 · mid-size D2C brands',
    },
    {
        'tag': 'Alert', 'tag_color': INK, 'tag_bg': RGBColor(0xED, 0xEF, 0xEE),
        'border': LINE,
        'title': 'High-Value Transaction Alert',
        'body': 'Transaction crosses a threshold you set. Slack ping to finance and Sheets row appended within 40 seconds. No morning report needed.',
        'stat': '40 seconds · threshold set by you',
        'stat_note': 'Transaction to Slack alert',
    },
]
aw = Inches(1.72)
agap = Inches(0.1)
ax_start = Inches(0.45)
ay = Inches(1.45)
ah = Inches(3.85)
tag_h = Inches(0.32)

for i, ag in enumerate(agents):
    ax = ax_start + i * (aw + agap)
    # header strip
    box(s, ax, ay, aw, tag_h, fill=ag['tag_bg'], line_color=ag['border'], line_w=Pt(1.2))
    tb(s, ag['tag'].upper(), ax + Inches(0.1), ay + Pt(5), aw * 0.55, Inches(0.25),
       size=Pt(8), bold=True, color=ag['tag_color'])
    if 'note' in ag:
        tb(s, ag['note'], ax + aw * 0.55, ay + Pt(7), aw * 0.45 - Inches(0.05), Inches(0.22),
           size=Pt(7), bold=True, color=ag['note_color'], align=PP_ALIGN.RIGHT)
    # body
    body_h = ah - tag_h - Inches(0.65)
    box(s, ax, ay + tag_h, aw, body_h, fill=PAPER, line_color=ag['border'], line_w=Pt(0.75))
    tb(s, ag['title'], ax + Inches(0.1), ay + tag_h + Pt(6), aw - Inches(0.15), Inches(0.45),
       size=Pt(9.5), bold=True, color=INK, wrap=True)
    tb(s, ag['body'], ax + Inches(0.1), ay + tag_h + Inches(0.52),
       aw - Inches(0.15), body_h - Inches(0.58),
       size=Pt(8.5), color=INK_MID, wrap=True)
    # stat footer
    ft_h = Inches(0.65)
    box(s, ax, ay + tag_h + body_h, aw, ft_h, fill=PAPER, line_color=ag['border'], line_w=Pt(0.75))
    box(s, ax, ay + tag_h + body_h, aw, Emu(9525), fill=ag['border'])
    tb(s, ag['stat'], ax + Inches(0.1), ay + tag_h + body_h + Pt(5), aw - Inches(0.15), Inches(0.3),
       size=Pt(9), bold=True, color=GREEN_INK)
    tb(s, ag['stat_note'], ax + Inches(0.1), ay + tag_h + body_h + Inches(0.32),
       aw - Inches(0.15), Inches(0.28), size=Pt(7.5), color=MUTED)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 08 · COD BEFORE / AFTER
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, PAPER)
chrome_bar(s, '08 / 10')

eyebrow_label(s, 'The most urgent problem first', Inches(0.55), Inches(0.5), Inches(5))

hl = s.shapes.add_textbox(Inches(0.55), Inches(0.8), Inches(9), Inches(0.65))
p = hl.text_frame.paragraphs[0]
r1 = p.add_run(); r1.text = 'Most RTO tools score the risk and stop. '
r1.font.name = FONT; r1.font.size = Pt(20); r1.font.color.rgb = INK
r2 = p.add_run(); r2.text = 'This agent confirms the order before dispatch.'
r2.font.name = FONT; r2.font.size = Pt(20); r2.font.color.rgb = INK; r2.font.italic = True

# Before panel (left)
bx = Inches(0.45); by = Inches(1.6); bw = Inches(4.5); bh = Inches(3.1)
box(s, bx, by, bw, bh, fill=PAPER, line_color=LINE, line_w=Pt(1))
tb(s, 'BEFORE · EVERY MORNING', bx + Inches(0.15), by + Pt(8), bw - Inches(0.3), Inches(0.25),
   size=Pt(7.5), bold=True, color=MUTED)

before_items = [
    'Open Cashfree. Export COD orders from last night.',
    'Open WhatsApp. Start sending confirmation messages.',
    "7 people didn't reply. Note them down. Send nudges at 2pm.",
    'Flag 3 high-risk orders to ops before they go for dispatch.',
    'Repeat tomorrow. And the day after.',
]
for idx, item in enumerate(before_items):
    iy = by + Inches(0.42) + idx * Inches(0.42)
    tb(s, '✕', bx + Inches(0.15), iy, Inches(0.25), Inches(0.35), size=Pt(11), bold=True, color=RED)
    tb(s, item, bx + Inches(0.42), iy, bw - Inches(0.55), Inches(0.38),
       size=Pt(9.5), bold=True, color=INK, wrap=True)

box(s, bx, by + bh - Inches(0.35), bw, Emu(9525), fill=AMBER)
tb(s, '1.5 to 2 hours, every single morning.',
   bx + Inches(0.15), by + bh - Inches(0.3), bw - Inches(0.2), Inches(0.28),
   size=Pt(10), bold=True, color=AMBER)

# Stats
sx = Inches(0.45); sy = by + bh + Inches(0.12)
sh = Inches(0.65); sw = Inches(2.2)
box(s, sx, sy, sw, sh, fill=PAPER, line_color=LINE)
tb(s, '26%', sx + Inches(0.12), sy + Pt(4), sw - Inches(0.2), Inches(0.38), size=Pt(28), color=INK)
tb(s, 'National RTO rate. Shipway 2025.', sx + Inches(0.12), sy + Inches(0.42),
   sw - Inches(0.2), Inches(0.2), size=Pt(8), color=INK_MID)
box(s, sx + sw + Inches(0.1), sy, sw, sh, fill=PAPER, line_color=LINE)
tb(s, '₹800', sx + sw + Inches(0.22), sy + Pt(4), sw - Inches(0.2), Inches(0.38), size=Pt(28), color=INK)
tb(s, 'Avg. loss per RTO. WareIQ.', sx + sw + Inches(0.22), sy + Inches(0.42),
   sw - Inches(0.2), Inches(0.2), size=Pt(8), color=INK_MID)

# After panel (right)
ax2 = Inches(5.2); ay2 = by; aw2 = Inches(4.6); ah2 = bh
box(s, ax2, ay2, aw2, ah2, fill=GREEN_TINT, line_color=GREEN, line_w=Pt(1.5))
tb(s, 'AFTER · COD CONFIRMATION AI AGENT', ax2 + Inches(0.15), ay2 + Pt(8),
   aw2 - Inches(0.3), Inches(0.25), size=Pt(7.5), bold=True, color=GREEN_INK)

after_items = [
    'COD order placed. Agent sends confirmation in 5 minutes.',
    'No reply in 4 hours. Agent sends nudge. Automatically.',
    'Still no reply. Ops flagged before dispatch. Order held.',
    'Runs at 2am on Sunday. Runs during your supplier meeting.',
    "You don't send a single message. You never open Cashfree for this again.",
]
for idx, item in enumerate(after_items):
    iy = ay2 + Inches(0.42) + idx * Inches(0.42)
    tb(s, '✓', ax2 + Inches(0.15), iy, Inches(0.25), Inches(0.35), size=Pt(11), bold=True, color=GREEN)
    tb(s, item, ax2 + Inches(0.42), iy, aw2 - Inches(0.55), Inches(0.38),
       size=Pt(9.5), bold=True, color=INK, wrap=True)

box(s, ax2, ay2 + ah2 - Inches(0.55), aw2, Emu(9525), fill=GREEN)
stat_tb = s.shapes.add_textbox(ax2 + Inches(0.15), ay2 + ah2 - Inches(0.5), aw2 - Inches(0.2), Inches(0.45))
p = stat_tb.text_frame.paragraphs[0]
r1 = p.add_run(); r1.text = '28–32%  →  16–20%'
r1.font.name = FONT; r1.font.size = Pt(22); r1.font.color.rgb = INK
tb(s, 'RTO rate within 30 days of activation. CampaignHQ 2026.',
   ax2 + Inches(0.15), ay2 + ah2 - Inches(0.22), aw2 - Inches(0.2), Inches(0.2),
   size=Pt(8), color=INK_MID)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 09 · WHAT RELAY DOES (4-row property table)
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, PAPER)
chrome_bar(s, '09 / 10')

eyebrow_label(s, 'Purpose-built for payment events', Inches(0.55), Inches(0.5), Inches(6))

hl = s.shapes.add_textbox(Inches(0.55), Inches(0.8), Inches(9), Inches(0.65))
p = hl.text_frame.paragraphs[0]
r1 = p.add_run(); r1.text = "Four things you can't get from a generic automation tool. "
r1.font.name = FONT; r1.font.size = Pt(20); r1.font.color.rgb = INK
r2 = p.add_run(); r2.text = "Because this isn't one."
r2.font.name = FONT; r2.font.size = Pt(20); r2.font.color.rgb = INK; r2.font.italic = True

prop_rows = [
    ('Two trigger types', 'Event or schedule. Both built in.',
     'Trigger on a Cashfree payment event (COD placed, payment failed, refund initiated, dispute raised) and the AI agent acts within minutes. Or run on a schedule: daily recon, weekly reports, recurring recovery sweeps. You pick which fits the job.',
     False),
    ('No developer needed', 'Activate from your dashboard.',
     'No webhook to parse. No API mapping. No logic to write before the job starts. Set your guardrails once: what to send, when to escalate, what threshold. The AI agent follows them on every qualifying event.',
     False),
    ('Data stays in Cashfree', 'PCI DSS Level 1 · SOC.',
     'Relay runs inside Cashfree\'s payment infrastructure. Your transaction data never has to move to a third-party server. Only the specific fields you choose reach WhatsApp or Slack. Everything else stays.',
     True),   # highlighted row
    ('Payment context built in', 'Not a generic workflow engine.',
     "The AI agents know what COD confirmation means. What a failed payment recovery looks like. What to flag before dispatch. What a settlement discrepancy is. You don't teach it payments from scratch. It already speaks the language.",
     False),
]

tx = Inches(0.45); ty = Inches(1.6)
tw2 = Inches(9.1); rh2 = Inches(0.84)
label_col = Inches(2.0)

box(s, tx, ty, tw2, len(prop_rows) * rh2, fill=None, line_color=INK, line_w=Pt(1.2))
for i, (title, sub, body, hilight) in enumerate(prop_rows):
    ry = ty + i * rh2
    row_fill = GREEN_TINT if hilight else (PAPER if i % 2 == 0 else RGBColor(0xF5, 0xF7, 0xF6))
    box(s, tx, ry, tw2, rh2, fill=row_fill, line_color=LINE, line_w=Pt(0.5))
    box(s, tx, ry, label_col, rh2, fill=RGBColor(0xEF, 0xF1, 0xF0) if not hilight else RGBColor(0xD8, 0xF3, 0xEA))
    box(s, tx + label_col, ry, Emu(9525), rh2, fill=LINE)
    tc = GREEN_INK if hilight else GREEN_INK
    tb(s, title, tx + Inches(0.12), ry + Pt(8), label_col - Inches(0.15), Inches(0.3),
       size=Pt(10), bold=True, color=tc)
    tb(s, sub, tx + Inches(0.12), ry + Inches(0.38), label_col - Inches(0.15), Inches(0.4),
       size=Pt(8.5), color=MUTED, wrap=True)
    tb(s, body, tx + label_col + Inches(0.15), ry + Pt(8),
       tw2 - label_col - Inches(0.2), rh2 - Pt(12),
       size=Pt(9.5), color=INK if hilight else INK_MID, wrap=True, bold=hilight)

# insight
iy = ty + len(prop_rows) * rh2 + Inches(0.12)
box(s, tx, iy, tw2, Inches(0.38), fill=RGBColor(0xF0, 0xF9, 0xF5))
tb(s, '=> The event is already inside Cashfree. The action runs from inside Cashfree. You just decide what happens in between.',
   tx + Inches(0.12), iy + Pt(4), tw2 - Inches(0.2), Inches(0.32),
   size=Pt(9.5), italic=True, color=INK_MID)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 · GET STARTED
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, PAPER)
chrome_bar(s, '10 / 10')

eyebrow_label(s, '15 minutes to your first live AI agent', Inches(0.55), Inches(0.5), Inches(6))

hl = s.shapes.add_textbox(Inches(0.55), Inches(0.82), Inches(7), Inches(0.5))
p = hl.text_frame.paragraphs[0]
r1 = p.add_run(); r1.text = 'Three steps. '
r1.font.name = FONT; r1.font.size = Pt(22); r1.font.color.rgb = INK
r2 = p.add_run(); r2.text = 'Then you do nothing.'
r2.font.name = FONT; r2.font.size = Pt(22); r2.font.italic = True; r2.font.color.rgb = INK

steps = [
    ('01', GREEN, GREEN_INK, GREEN_TINT, 'Pick an AI agent',
     "Go to Relay inside your Cashfree dashboard. Browse the AI agents. Pick the one for the problem you feel most right now: COD confirmations eating your morning, failed payments you're chasing manually, or settlement recon taking your Fridays."),
    ('02', MUTED, MUTED, RGBColor(0xF0, 0xF2, 0xF1), 'Set your guardrails',
     'What message to send. When to send a nudge. When to flag ops. What threshold to use. 10 minutes. The AI agent follows these rules on every qualifying event from that point forward. No further input from you.'),
    ('03', GREEN, GREEN_INK, GREEN_TINT, 'You do nothing',
     'The AI agent runs on every event. Day and night. Weekend. Holiday. Supplier meeting. Your only job from here: decide which AI agent to activate next. The morning you\'ve been doing manually every day stops being yours.'),
]
sw2 = Inches(2.85)
sgap = Inches(0.2)
sx_start = Inches(0.45)
sy = Inches(1.5)
sh2 = Inches(2.0)

for i, (num, num_c, lbl_c, hdr_bg, lbl, body) in enumerate(steps):
    sx2 = sx_start + i * (sw2 + sgap)
    hh = Inches(0.42)
    box(s, sx2, sy, sw2, hh, fill=hdr_bg, line_color=num_c, line_w=Pt(1.2))
    tb(s, num, sx2 + Inches(0.12), sy + Pt(3), Inches(0.5), hh, size=Pt(22), bold=False, color=num_c)
    tb(s, lbl.upper(), sx2 + Inches(0.6), sy + Pt(10), sw2 - Inches(0.7), hh,
       size=Pt(7.5), bold=True, color=lbl_c)
    box(s, sx2, sy + hh, sw2, sh2 - hh, fill=PAPER, line_color=num_c if num_c != MUTED else LINE, line_w=Pt(0.75))
    tb(s, body, sx2 + Inches(0.12), sy + hh + Pt(8), sw2 - Inches(0.2), sh2 - hh - Pt(12),
       size=Pt(9.5), color=INK_MID, wrap=True)

# Property box
px = Inches(0.45); py = sy + sh2 + Inches(0.2); pw = Inches(9.1)
props = [
    ('Pricing', 'Free to activate during Early Beta. Per-run pricing at GA: directionally ₹0.50 to ₹1 per run. WhatsApp template costs at cost. No platform fee during Early Beta.', False),
    ('Data',    'PCI DSS Level 1 and SOC-certified. Same infrastructure that processes your payments. Your transaction data does not leave Cashfree.', False),
    ('Access',  'Ask for access on your next Cashfree call, or write to mothi.venkatesh@cashfree.com. We\'re activating merchants in batches during Early Beta. Currently free.', True),
]
lbl_w = Inches(1.1)
prh = Inches(0.55)
box(s, px, py, pw, len(props) * prh, fill=None, line_color=INK, line_w=Pt(1.2))
for i, (plbl, pval, hilight) in enumerate(props):
    pry = py + i * prh
    row_fill = GREEN_TINT if hilight else PAPER
    box(s, px, pry, pw, prh, fill=row_fill, line_color=LINE, line_w=Pt(0.5))
    box(s, px, pry, lbl_w, prh, fill=RGBColor(0xEF, 0xF1, 0xF0) if not hilight else RGBColor(0xD8, 0xF3, 0xEA))
    box(s, px + lbl_w, pry, Emu(9525), prh, fill=LINE)
    tb(s, plbl.upper(), px + Inches(0.1), pry + Pt(6), lbl_w - Inches(0.12), prh,
       size=Pt(8), bold=True, color=GREEN_INK if hilight else MUTED)
    tb(s, pval, px + lbl_w + Inches(0.12), pry + Pt(6), pw - lbl_w - Inches(0.2), prh - Pt(8),
       size=Pt(9), color=INK if hilight else INK_MID, wrap=True, bold=hilight)


# ═══════════════════════════════════════════════════════════════════════
# CLOSING · DARK
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, DARK)
chrome_bar(s, '··', dark=True)

# Large RELAY watermark text
for i, (font_size, y_frac, opacity_val) in enumerate([
    (200, 0.45, 0.06), (148, 0.85, 0.04), (110, 1.15, 0.03)
]):
    wm = s.shapes.add_textbox(Inches(-0.5), H * y_frac - Inches(1.5), W + Inches(1), Inches(2))
    wm.text_frame.word_wrap = False
    p = wm.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = 'RELAY'
    r.font.name = FONT; r.font.size = Pt(font_size); r.font.bold = True
    grey_val = int(255 * opacity_val * 3)
    r.font.color.rgb = RGBColor(grey_val, grey_val + 4, grey_val + 2)

eyebrow_label(s, 'The only question left', Inches(0.55), Inches(1.5),
              Inches(5), color=GREEN_SOFT)

hl = s.shapes.add_textbox(Inches(0.55), Inches(1.9), Inches(9), Inches(2.2))
hl.text_frame.word_wrap = True
p1 = hl.text_frame.paragraphs[0]
r1 = p1.add_run(); r1.text = 'Which problem do you want to'
r1.font.name = FONT; r1.font.size = Pt(44); r1.font.color.rgb = OFF_WHITE
p2 = hl.text_frame.add_paragraph()
r2 = p2.add_run(); r2.text = 'stop doing by hand?'
r2.font.name = FONT; r2.font.size = Pt(44); r2.font.color.rgb = GREEN_SOFT; r2.font.italic = True

# CTA box
box(s, Inches(0.45), Inches(4.15), Inches(5.8), Inches(0.88),
    fill=RGBColor(0x16, 0x2C, 0x26), line_color=GREEN, line_w=Pt(1.5))
box(s, Inches(0.45), Inches(4.15), Emu(12700), Inches(0.88), fill=GREEN)
tb(s, 'Write to mothi.venkatesh@cashfree.com to get activated during Early Beta.',
   Inches(0.62), Inches(4.3), Inches(5.5), Inches(0.55),
   size=Pt(11), color=OFF_WHITE, wrap=True)
# highlight email
tb(s, 'mothi.venkatesh@cashfree.com',
   Inches(0.62), Inches(4.52), Inches(5.5), Inches(0.28),
   size=Pt(10.5), bold=True, color=GREEN_SOFT)


# ── save ──────────────────────────────────────────────────────────────
out = '/Users/mothi.venkatesh/Documents/Cashfree_Relay_Pitch.pptx'
prs.save(out)
print(f'Saved: {out}')
